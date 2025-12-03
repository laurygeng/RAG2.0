#!/usr/bin/env python3
"""
generate_pptx.py

Create a PPTX summarizing RAG system architecture, metrics and include diagnostics samples.
Reads the latest summary and diagnostics files under evaluation/ and the augmented results CSV.
Writes a pptx to results/ with timestamp.
"""
import os
import sys
import json
from datetime import datetime
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
EVAL_DIR = os.path.join(ROOT, 'evaluation')
RESULTS_DIR = os.path.join(ROOT, 'results')
os.makedirs(RESULTS_DIR, exist_ok=True)

# find latest evaluation folder
eval_folders = []
if os.path.exists(EVAL_DIR):
    eval_folders = [os.path.join(EVAL_DIR, d) for d in os.listdir(EVAL_DIR) if os.path.isdir(os.path.join(EVAL_DIR, d))]
if not eval_folders:
    print('No evaluation folders found under', EVAL_DIR)
    sys.exit(1)
latest_eval = max(eval_folders, key=os.path.getmtime)

# paths
summary_csv = os.path.join(latest_eval, [f for f in os.listdir(latest_eval) if f.startswith('real_summary_')][0])
llama_eval_csv = os.path.join(latest_eval, [f for f in os.listdir(latest_eval) if f.startswith('real_llama_evaluation_')][0])
low_perf_csv_candidates = [f for f in os.listdir(latest_eval) if f.startswith('low_performance_')]
low_perf_csv = os.path.join(latest_eval, low_perf_csv_candidates[0]) if low_perf_csv_candidates else None

# results csv (use latest in results/ matching real_answers_with_retrieval)
res_root = os.path.join(ROOT, 'results')
res_files = [os.path.join(res_root, f) for f in os.listdir(res_root) if f.startswith('real_answers_with_retrieval_')]
if not res_files:
    print('No augmented results CSV found under', res_root)
    sys.exit(1)
results_csv = max(res_files, key=os.path.getmtime)

summary_df = pd.read_csv(summary_csv)
llama_df = pd.read_csv(llama_eval_csv)
results_df = pd.read_csv(results_csv)
low_df = pd.read_csv(low_perf_csv) if low_perf_csv else pd.DataFrame()

# Build presentation
prs = Presentation()
# helper
def add_title_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
    return slide

def add_table_slide(title, df, max_rows=10):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows = min(len(df)+1, max_rows+1)
    cols = len(df.columns)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.2*rows)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # header
    for j, c in enumerate(df.columns):
        table.cell(0, j).text = str(c)
    for i in range(1, rows):
        if i-1 >= len(df):
            for j in range(cols):
                table.cell(i, j).text = ''
            continue
        for j, c in enumerate(df.columns):
            val = df.iloc[i-1][c]
            cell = table.cell(i, j)
            cell.text = str(val)
    return slide

# Slides content
add_title_slide('RAG Architecture & Evaluation Report', 'Based on local Ollama + Chroma — Summary and diagnostic samples')
add_bullet_slide('Agenda', [
    'System architecture overview', 'Metric definitions', 'Key experimental results', 'Threshold & reranking sensitivity', 'Diagnostic examples', 'Conclusions & recommendations'
])

# Architecture slide (textual)
add_bullet_slide('System Architecture (Brief)', [
    'Data Ingestion -> Chunking -> Vector DB (Chroma)',
    'Retriever (all-MiniLM-L6-v2) -> Top-K candidates',
    'Optional Reranker -> Evidence-first Prompt -> Local LLM (Ollama)',
    'Results CSV -> Evaluator (evaluate_answers.py)'
])

# Metrics slide: include summary table (first row)
add_table_slide('Key Metrics (Summary)', summary_df.round(4))

# Add a slide with sensitivity bullet points
add_bullet_slide('Threshold & decision sensitivity (summary)', [
    'Lowering support_threshold (0.75 -> 0.6) can notably increase Faithfulness',
    'Thresholds for token-overlap are very sensitive for Relevant_Retrieved_Proportion (e.g., 0.1 -> 0.99; 0.25 -> 0.35)',
    'Semantic decision (cosine >= 0.6) detects more semantically supported spans'
])

# Diagnostics slides: include top 3 low-performance samples (or up to 4)
num_examples = min(4, len(low_df))
if num_examples == 0:
    add_bullet_slide('Diagnostic examples', ['No low-performance rows found in diagnostics.'])
else:
    for i in range(num_examples):
        row = low_df.iloc[i]
        q = row.get('Question', '')
        gt = row.get('Ground_Truth', '')
        ans = row.get('Real_LLaMA_Answer', '')
        # get retrieved passages from results_df matching Index or Question
        idx = row.get('Index', None)
        retrieved = ''
        try:
            match = results_df[results_df.index == idx]
            if match.empty:
                match = results_df[results_df['Question'] == q]
            if not match.empty:
                rp = match.iloc[0].get('Retrieved_Passages', '')
                try:
                    docs = json.loads(rp) if isinstance(rp, str) else []
                except Exception:
                    docs = str(rp).split('||') if rp else []
                retrieved = '\n\n'.join(docs[:5])
        except Exception:
            retrieved = ''
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f'Diagnostic Sample {i+1} (Index {idx})'
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.text = f'Q: {q}\n\nGT: {gt}\n\nAnswer: {ans}\n\nTop retrieved passages:\n{retrieved}'

# Appendix: commands
add_bullet_slide('Reproduction commands (partial)', [
    'python3 code/evaluate_answers.py /path/to/results_with_retrieval.csv',
    'EVAL_RELEVANCE_METHOD=semantic EVAL_SEMANTIC_RELEVANCE_THRESHOLD=0.6 \\',
    'EVAL_SUPPORT_THRESHOLD=0.6 python3 code/evaluate_answers.py /path/to/results_with_retrieval.csv',
    'python3 code/rerank_experiment.py /path/to/results_with_retrieval.csv 20 --method semantic --sem_thresh 0.6'
])

# Save PPTX
ts = datetime.now().strftime('%Y%m%d_%H%M%S')
out_path = os.path.join(RESULTS_DIR, f'RAG_evaluation_presentation_{ts}.pptx')
prs.save(out_path)
print('Wrote PPTX to', out_path)
