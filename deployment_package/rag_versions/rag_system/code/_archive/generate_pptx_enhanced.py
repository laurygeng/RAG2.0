#!/usr/bin/env python3
"""
generate_pptx_enhanced.py

Enhanced PPTX generator:
- Creates a simple architecture diagram image (PIL)
- Computes sensitivity curves for token_overlap and semantic thresholds
  by evaluating the latest augmented results CSV
- Generates plots (matplotlib)
- Builds an English PPTX with nicer tables, images and plots

Writes: results/RAG_evaluation_presentation_enhanced_<ts>.pptx
"""
import os
import sys
from datetime import datetime
import json

import pandas as pd
import numpy as np

# plotting and image
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

# pptx
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
EVAL_DIR = os.path.join(ROOT, 'evaluation')
RESULTS_DIR = os.path.join(ROOT, 'results')
os.makedirs(RESULTS_DIR, exist_ok=True)

# helper: latest files
def latest_eval_folder():
    if not os.path.exists(EVAL_DIR):
        return None
    folders = [os.path.join(EVAL_DIR, d) for d in os.listdir(EVAL_DIR) if os.path.isdir(os.path.join(EVAL_DIR,d))]
    return max(folders, key=os.path.getmtime) if folders else None

le = latest_eval_folder()
if not le:
    print('No evaluation found under', EVAL_DIR)
    sys.exit(1)

summary_csv = [os.path.join(le,f) for f in os.listdir(le) if f.startswith('real_summary_')][0]
llama_csv = [os.path.join(le,f) for f in os.listdir(le) if f.startswith('real_llama_evaluation_')][0]
low_csvs = [os.path.join(le,f) for f in os.listdir(le) if f.startswith('low_performance_')]
low_csv = low_csvs[0] if low_csvs else None

# latest augmented results
res_root = os.path.join(ROOT, 'results')
res_files = [os.path.join(res_root, f) for f in os.listdir(res_root) if f.startswith('real_answers_with_retrieval_')]
if not res_files:
    print('No augmented results csv in', res_root)
    sys.exit(1)
results_csv = max(res_files, key=os.path.getmtime)

summary_df = pd.read_csv(summary_csv)
llama_df = pd.read_csv(llama_csv)
results_df = pd.read_csv(results_csv)
low_df = pd.read_csv(low_csv) if low_csv else pd.DataFrame()

# --- create architecture diagram (PIL) ---
ARCH_PNG = os.path.join(RESULTS_DIR, 'arch_diagram.png')
W, H = 1200, 600
img = Image.new('RGB', (W,H), color='white')
d = ImageDraw.Draw(img)

# simple boxes
boxes = [
    ('Data\nIngestion', 80, 80),
    ('Chunking\n& KB Pickle', 320, 80),
    ('Vector DB\n(Chroma)', 560, 80),
    ('Retriever\n(all-MiniLM-L6-v2)', 800, 80),
    ('Reranker\n(optional)', 320, 300),
    ('Prompt\n(evidence-first)', 560, 300),
    ('Local LLM\n(Ollama)', 800, 300),
    ('Results CSV\n& Evaluator', 960, 420),
]

for label, x, y in boxes:
    d.rectangle([x, y, x+220, y+100], outline='black', width=2)
    # text center
    lines = label.split('\n')
    fy = y+12
    for line in lines:
        d.text((x+10, fy), line, fill='black')
        fy += 18

# arrows
def arrow(a,b):
    d.line([a,b], fill='black', width=2)

arrow((300,130),(320,130))
arrow((540,130),(560,130))
arrow((760,130),(800,130))
arrow((920,350),(960,350))
arrow((660,200),(660,300))
arrow((900,200),(900,300))

img.save(ARCH_PNG)
print('Wrote architecture diagram to', ARCH_PNG)

# --- sensitivity computation ---
# We'll compute token-overlap sweep and, if embedding model available, semantic sweep
from math import isfinite
try:
    from sentence_transformers import SentenceTransformer
    EMB = SentenceTransformer('all-MiniLM-L6-v2')
    has_emb = True
except Exception:
    EMB = None
    has_emb = False

# utility tokenization
import re

def tokenize(text):
    if pd.isna(text) or not text:
        return []
    s = re.sub(r"[^0-9a-zA-Z]+", " ", str(text).lower())
    toks = [t for t in s.split() if t]
    return toks

# is_relevant for token overlap
def is_relevant_token(doc, gt, thresh):
    gtoks = set(tokenize(gt))
    if len(gtoks) == 0:
        return False
    dtoks = set(tokenize(doc))
    return (len(dtoks & gtoks) / len(gtoks)) >= thresh

# semantic relevance
from sklearn.metrics.pairwise import cosine_similarity

def is_relevant_sem(doc, gt, thresh):
    if not has_emb:
        return False
    try:
        ge = EMB.encode(gt)
        de = EMB.encode(doc)
        sim = float(cosine_similarity([ge],[de])[0][0])
        return sim >= thresh
    except Exception:
        return False

# compute metric for a given method and threshold
def compute_metrics(method='token', thresh=0.5, top_k=20, recall_k=5):
    props = []
    recalls = []
    for idx, row in results_df.iterrows():
        rp = row.get('Retrieved_Passages', '')
        try:
            docs = json.loads(rp) if isinstance(rp, str) else []
        except Exception:
            docs = str(rp).split('||') if rp else []
        topk = docs[:top_k]
        gt = row.get('Ground_Truth_Answer', '')
        if method == 'token':
            rels = [1 if is_relevant_token(d, gt, thresh) else 0 for d in topk]
            props.append(sum(rels)/len(topk) if topk else 0)
            recalls.append(sum(rels[:recall_k]))
        elif method == 'semantic':
            rels = [1 if is_relevant_sem(d, gt, thresh) else 0 for d in topk]
            props.append(sum(rels)/len(topk) if topk else 0)
            recalls.append(sum(rels[:recall_k]))
    return float(np.mean(props)), float(np.mean(recalls))

# sweeps
token_threshs = np.linspace(0.1, 0.5, 9)
sem_threshs = np.linspace(0.5, 0.8, 7) if has_emb else []

token_props = []
token_recalls = []
for t in token_threshs:
    p, r = compute_metrics('token', float(t))
    token_props.append(p)
    token_recalls.append(r)

sem_props = []
sem_recalls = []
for s in sem_threshs:
    p, r = compute_metrics('semantic', float(s))
    sem_props.append(p)
    sem_recalls.append(r)

# plot token overlap sensitivity
TOK_PNG = os.path.join(RESULTS_DIR, 'sensitivity_token.png')
plt.figure(figsize=(6,4))
plt.plot(token_threshs, token_props, marker='o', label='Avg Relevant Proportion')
plt.plot(token_threshs, token_recalls, marker='x', label='Avg Recall@5 (count)')
plt.xlabel('Token-overlap threshold')
plt.ylabel('Metric')
plt.title('Token-overlap Sensitivity')
plt.legend()
plt.grid(True)
plt.tight_layout()
plt.savefig(TOK_PNG)
plt.close()
print('Wrote token sensitivity plot to', TOK_PNG)

# plot semantic sensitivity if available
SEM_PNG = None
if has_emb:
    SEM_PNG = os.path.join(RESULTS_DIR, 'sensitivity_semantic.png')
    plt.figure(figsize=(6,4))
    plt.plot(sem_threshs, sem_props, marker='o', label='Avg Relevant Proportion')
    plt.plot(sem_threshs, sem_recalls, marker='x', label='Avg Recall@5 (count)')
    plt.xlabel('Semantic cosine threshold')
    plt.ylabel('Metric')
    plt.title('Semantic Threshold Sensitivity')
    plt.legend()
    plt.grid(True)
    plt.tight_layout()
    plt.savefig(SEM_PNG)
    plt.close()
    print('Wrote semantic sensitivity plot to', SEM_PNG)

# --- Build PPTX (English) ---
prs = Presentation()
# title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'RAG System Architecture & Evaluation'
slide.placeholders[1].text = 'Local Ollama + Chroma — Enhanced Summary'

# agenda
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Agenda'
body = slide.shapes.placeholders[1].text_frame
body.clear()
for b in ['System architecture', 'Metrics (definitions)', 'Key results', 'Sensitivity analysis', 'Diagnostics', 'Conclusions & next steps']:
    p = body.add_paragraph()
    p.text = b

# architecture slide with image
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Architecture (overview)'
left = Inches(0.5)
top = Inches(1.5)
pic = slide.shapes.add_picture(ARCH_PNG, left, top, width=Inches(9))

# summary table (english)
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Key Metrics (summary)'
# create smaller table with selected columns
cols = ['Semantic_Similarity','Jaccard_Similarity','Overlap_Similarity','Combined_Score','Faithfulness','Precision','Recall','F1','Recall_at_5','Relevant_Retrieved_Proportion','Retrieved_Context_Score','Supported_Sentences','Accuracy']
rowvals = [summary_df.iloc[0].get(c, '') for c in cols]

rows = 2
cols_n = len(cols)
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(1.2)
table = slide.shapes.add_table(rows, cols_n, left, top, width, height).table
for j,c in enumerate(cols):
    table.cell(0,j).text = c
    table.cell(1,j).text = str(rowvals[j])

# sensitivity slide insert token plot
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Token-overlap Sensitivity'
slide.shapes.add_picture(TOK_PNG, Inches(1), Inches(1.6), width=Inches(8))

if SEM_PNG:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Semantic Threshold Sensitivity'
    slide.shapes.add_picture(SEM_PNG, Inches(1), Inches(1.6), width=Inches(8))

# diagnostics: include up to 4 examples with retrieved passages
num_examples = min(4, len(low_df))
if num_examples==0:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Diagnostics'
    slide.shapes.placeholders[1].text = 'No low-performance diagnostics found.'
else:
    for i in range(num_examples):
        row = low_df.iloc[i]
        idx = row.get('Index')
        q = row.get('Question','')
        gt = row.get('Ground_Truth','')
        ans = row.get('Real_LLaMA_Answer','')
        # find in results_df by index
        match = results_df[results_df.index==idx]
        if match.empty:
            match = results_df[results_df['Question']==q]
        retrieved = ''
        if not match.empty:
            rp = match.iloc[0].get('Retrieved_Passages','')
            try:
                docs = json.loads(rp) if isinstance(rp,str) else []
            except Exception:
                docs = str(rp).split('||') if rp else []
            retrieved = '\n\n'.join(docs[:5])
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f'Diagnostic example {i+1} (Index {idx})'
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        tf.text = f'Q: {q}\n\nGT: {gt}\n\nAnswer: {ans}\n\nTop retrieved passages:\n{retrieved}'

# conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Conclusions & Next Steps'
body = slide.shapes.placeholders[1].text_frame
body.clear()
for p in [
    'Semantic similarity is high but token-level overlap is low (rewrite/abstraction).',
    'Consider semantic relevance or hybrid rules for Relevant_Retrieved metric.',
    'Short-term: lower support_threshold and use semantic checks; mid-term: reranker + hybrid retrieval.'
]:
    pp = body.add_paragraph()
    pp.text = p

# save PPTX
ts = datetime.now().strftime('%Y%m%d_%H%M%S')
out = os.path.join(RESULTS_DIR, f'RAG_evaluation_presentation_enhanced_{ts}.pptx')
prs.save(out)
print('Wrote enhanced PPTX to', out)

if has_emb:
    print('Embedding model used for semantic sensitivity (all-MiniLM-L6-v2).')
else:
    print('Embedding model NOT available; semantic sensitivity skipped.')
